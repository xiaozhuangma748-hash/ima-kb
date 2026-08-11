"""多格式文档解析器。

支持格式：
- PDF (.pdf)            → PyMuPDF（扫描版自动走 OCR）
- Word (.docx)          → python-docx
- Word (.doc)           → macOS textutil
- Excel (.xlsx)         → openpyxl
- PowerPoint (.pptx)    → python-pptx
- 图片 (.png/.jpg/...)  → PaddleOCR（优先）/ Tesseract（降级）
- Markdown (.md/.markdown) → 纯文本读取
- 文本 (.txt/.log)      → 纯文本读取
- 代码 (.py/.js/...)    → 纯文本读取，带语言标签
- HTML (.html/.htm)     → trafilatura 抽取正文

OCR 依赖（可选，任一即可）：
- PaddleOCR（推荐，精度高）：pip install paddlepaddle paddleocr
- Tesseract（降级）：brew install tesseract tesseract-lang + pip install pytesseract pillow

OCR 流程：图片预处理（灰度+二值化+放大）→ PaddleOCR → 降级 Tesseract

返回统一的 ParsedDocument 结构。
"""
from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Dict, Optional


@dataclass
class ParsedDocument:
    """解析后的统一文档结构。

    Attributes:
        text: 提取出的纯文本内容
        title: 文档标题（默认为文件名，不带扩展名）
        file_path: 原文件路径
        file_type: 文件类型（扩展名，小写）
        language: 内容语言（如 'zh' / 'en' / 'code' / 'unknown'）
        meta: 额外元信息（页数、作者等）
        format_tag: 格式标签（'excel'/'ppt'/'code'/'markdown'/'docx'/'pdf'/'image'/'text'/'html'），
                    用于下游 chunker 和 chain 做格式感知处理
    """

    text: str
    title: str
    file_path: Path
    file_type: str
    language: str = "unknown"
    meta: Dict[str, str] = field(default_factory=dict)
    format_tag: str = "text"


# ---- 支持的文件类型 ----
# 注：代码文件（.py/.js/.ts/...）按用户决策不入库，已从支持列表移除
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".pptx",
    ".md", ".markdown",
    ".txt", ".log",
    ".html", ".htm",
    # 图片（OCR）
    ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp",
}

# 代码文件后缀 → 语言标签
_CODE_LANGUAGES: Dict[str, str] = {
    ".py": "python", ".js": "javascript", ".ts": "typescript",
    ".tsx": "typescript", ".jsx": "javascript",
    ".java": "java", ".go": "go", ".rs": "rust",
    ".c": "c", ".cpp": "cpp", ".h": "c", ".hpp": "cpp",
    ".cs": "csharp", ".rb": "ruby", ".php": "php",
    ".swift": "swift", ".kt": "kotlin", ".scala": "scala",
    ".sh": "shell", ".bash": "shell",
    ".sql": "sql", ".yaml": "yaml", ".yml": "yaml",
    ".json": "json", ".xml": "xml", ".toml": "toml",
    ".ini": "ini", ".conf": "conf",
}


class ParseError(Exception):
    """文档解析失败。"""


# ============================================================
# OCR 工具（PaddleOCR 优先，Tesseract 降级，可选依赖）
# ============================================================

import threading as _threading

# OCR 引擎检测缓存
_ocr_checked = False
_ocr_available = False

# PaddleOCR 单例（初始化慢，全局复用）
_paddle_ocr = None
_paddle_ocr_failed = False  # 标记 PaddleOCR 不可用，避免反复尝试

# PaddleOCR 调用锁：predict() 不保证线程安全，并行 OCR 时串行化 PaddleOCR 调用
# Tesseract 走子进程，无需此锁
_paddle_ocr_lock = _threading.Lock()

# PaddleX layout_parsing 单例（版面分析，首次加载慢）
_layout_pipeline = None
_layout_pipeline_failed = False


def _get_layout_pipeline():
    """获取 PaddleX layout_parsing pipeline 单例（懒加载）。

    layout_parsing 集成了版面检测（17类区域）+ 表格识别（SLANet_plus）+
    印章识别 + 文档预处理（方向矫正+去弯曲），输出结构化的 parsing_res_list。

    Returns:
        pipeline 实例，或 None（不可用）
    """
    global _layout_pipeline, _layout_pipeline_failed
    if _layout_pipeline_failed:
        return None
    if _layout_pipeline is not None:
        return _layout_pipeline
    try:
        import logging
        import os as _os
        # 静默 PaddleX 日志
        for name in ("paddleocr", "paddle", "paddlex", "ppocr"):
            logging.getLogger(name).setLevel(logging.ERROR)
        # PaddleX 用 print 输出模型加载信息，重定向 stdout/stderr
        _devnull_fd = _os.open(_os.devnull, _os.O_WRONLY)
        _saved_stdout_fd = _os.dup(1)
        _saved_stderr_fd = _os.dup(2)
        _os.dup2(_devnull_fd, 1)
        _os.dup2(_devnull_fd, 2)
        try:
            from paddlex import create_pipeline  # type: ignore
            _layout_pipeline = create_pipeline(pipeline="layout_parsing", device="cpu")
        finally:
            _os.dup2(_saved_stdout_fd, 1)
            _os.dup2(_saved_stderr_fd, 2)
            _os.close(_devnull_fd)
            _os.close(_saved_stdout_fd)
            _os.close(_saved_stderr_fd)
        return _layout_pipeline
    except Exception:
        _layout_pipeline_failed = True
        return None


def _check_ocr() -> bool:
    """检测 OCR 是否可用（PaddleOCR 或 Tesseract 任一可用即可）。"""
    global _ocr_checked, _ocr_available
    if _ocr_checked:
        return _ocr_available
    _ocr_checked = True
    # 1. 检测 PaddleOCR
    try:
        from paddleocr import PaddleOCR  # type: ignore  # noqa: F401
        _ocr_available = True
        return True
    except ImportError:
        pass
    # 2. 降级：检测 Tesseract
    try:
        import pytesseract  # type: ignore  # noqa: F401
        from shutil import which
        if which("tesseract") is not None:
            _ocr_available = True
            return True
    except ImportError:
        pass
    _ocr_available = False
    return False


def reset_ocr_cache() -> None:
    """重置 OCR 可用性检测缓存。

    应用场景：用户在 REPL 运行期间安装了 OCR 依赖，
    调用此函数后下次解析会重新检测，无需重启进程。
    """
    global _ocr_checked, _ocr_available, _paddle_ocr, _paddle_ocr_failed
    global _layout_pipeline, _layout_pipeline_failed
    _ocr_checked = False
    _ocr_available = False
    _paddle_ocr = None
    _paddle_ocr_failed = False
    _layout_pipeline = None
    _layout_pipeline_failed = False


def _get_paddle_ocr():
    """获取 PaddleOCR 单例（懒加载）。

    Returns:
        PaddleOCR 实例，或 None（不可用）
    """
    global _paddle_ocr, _paddle_ocr_failed
    if _paddle_ocr_failed:
        return None
    if _paddle_ocr is not None:
        return _paddle_ocr
    try:
        import logging
        import os as _os
        import sys as _sys
        # 静默 PaddleOCR / PaddleX 的日志（print + logging 双管齐下）
        logging.getLogger("paddleocr").setLevel(logging.ERROR)
        logging.getLogger("paddle").setLevel(logging.ERROR)
        logging.getLogger("paddlex").setLevel(logging.ERROR)
        logging.getLogger("ppocr").setLevel(logging.ERROR)
        # PaddleOCR 用 print 输出模型加载信息，重定向 stdout/stderr
        _devnull_fd = _os.open(_os.devnull, _os.O_WRONLY)
        _saved_stdout_fd = _os.dup(1)
        _saved_stderr_fd = _os.dup(2)
        _os.dup2(_devnull_fd, 1)
        _os.dup2(_devnull_fd, 2)
        try:
            from paddleocr import PaddleOCR  # type: ignore
            _paddle_ocr = PaddleOCR(lang="ch")
        finally:
            _os.dup2(_saved_stdout_fd, 1)
            _os.dup2(_saved_stderr_fd, 2)
            _os.close(_devnull_fd)
            _os.close(_saved_stdout_fd)
            _os.close(_saved_stderr_fd)
        return _paddle_ocr
    except Exception:
        _paddle_ocr_failed = True
        return None


def _preprocess_image(image):
    """图片预处理：灰度化 + 二值化 + 放大，提升 OCR 识别率。

    Args:
        image: PIL.Image 对象

    Returns:
        预处理后的 PIL.Image（L 模式，灰度二值化）
    """
    from PIL import Image, ImageOps  # type: ignore

    # 1. 转灰度
    if image.mode != "L":
        image = image.convert("L")

    # 2. 如果分辨率较低，放大 2 倍（提升小字识别）
    w, h = image.size
    if w < 1000:
        image = image.resize((w * 2, h * 2), Image.LANCZOS)

    # 3. 自动对比度
    image = ImageOps.autocontrast(image)

    # 4. Otsu 自适应二值化（对扫描公文效果显著）
    import numpy as np  # type: ignore
    arr = np.array(image)
    # 简单 Otsu：用 PIL 内置的点操作近似
    threshold = arr.mean()
    arr_bin = (arr > threshold) * 255
    image = Image.fromarray(arr_bin.astype("uint8"), mode="L")

    return image


def _ocr_image_paddle(image) -> str:
    """用 PaddleOCR 识别图片，返回文本。

    线程安全：predict() 不保证可重入，用 _paddle_ocr_lock 串行化调用。
    并行 OCR 场景下，多线程会在此锁处排队，但 Tesseract 分支仍可并行。

    Args:
        image: PIL.Image 对象

    Returns:
        识别出的文本
    """
    import numpy as np  # type: ignore
    ocr = _get_paddle_ocr()
    if ocr is None:
        return ""
    # PIL → numpy array
    arr = np.array(image)
    # PaddleOCR.predict() 可能非线程安全，加锁串行化
    with _paddle_ocr_lock:
        result = ocr.predict(arr)
    if not result:
        return ""
    # PaddleOCR 3.x 返回结果：list of dict-like
    texts = []
    for item in result:
        # 3.x API：item.json['res']['rec_texts']
        if hasattr(item, "json"):
            res = item.json.get("res", {})
            rec_texts = res.get("rec_texts", [])
            texts.extend(rec_texts)
        elif isinstance(item, dict):
            rec_texts = item.get("rec_texts", [])
            texts.extend(rec_texts)
    return "\n".join(texts)


def _ocr_image_tesseract(image) -> str:
    """用 Tesseract 识别图片（降级方案）。"""
    import pytesseract  # type: ignore
    return pytesseract.image_to_string(image, lang="chi_sim+eng")


def _clean_ocr_text(text: str) -> str:
    """OCR 文本后处理：清理乱码、合并断行、去除中文间空格。

    处理：
    1. 去除 OCR 误识别的几何符号/特殊符号（■●□◆○等）
    2. 去除中文字符之间的空格（OCR 常随机插入）
    3. 合并视觉换行（复用 _merge_visual_line_breaks）
    4. 压缩连续空格
    """
    if not text or not text.strip():
        return text

    # 1. 去除 OCR 常见乱码符号（几何图形、装饰符号）
    # 保留正文标点和常规符号
    text = re.sub(r"[\u25A0-\u25FF\u2600-\u27BF\u2B00-\u2BFF]", "", text)

    # 2. 去除中文字符之间的空格
    # 模式：中文+空格+中文 → 中文+中文（只处理单空格，避免破坏英文短语）
    text = re.sub(r"([\u4e00-\u9fff])\s+([\u4e00-\u9fff])", r"\1\2", text)
    # 连续处理多次（处理"中 中 中"这种情况）
    for _ in range(3):
        new_text = re.sub(r"([\u4e00-\u9fff])\s+([\u4e00-\u9fff])", r"\1\2", text)
        if new_text == text:
            break
        text = new_text

    # 3. 压缩连续空格（3+ 空格 → 1 个）
    text = re.sub(r"[ \t]{3,}", " ", text)

    # 4. 合并视觉换行
    text = _merge_visual_line_breaks(text)

    return text


def _html_table_to_markdown(html: str) -> str:
    """把 PaddleX 表格识别输出的 HTML 转为 Markdown 表格。

    PaddleX table_recognition 输出 <table><tr><td>...</td></tr></table> 格式，
    转为 Markdown `| cell | cell |` 格式以提升检索友好度。
    解析失败时返回原始 HTML。
    """
    if not html or "<" not in html:
        return html
    try:
        from html.parser import HTMLParser

        class _TableParser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.rows: list[list[str]] = []
                self._cur_row: list[str] | None = None
                self._cur_cell: list[str] = []
                self._in_cell = False

            def handle_starttag(self, tag, attrs):
                if tag == "tr":
                    self._cur_row = []
                elif tag in ("td", "th"):
                    self._in_cell = True
                    self._cur_cell = []

            def handle_endtag(self, tag):
                if tag == "tr" and self._cur_row is not None:
                    self.rows.append(self._cur_row)
                    self._cur_row = None
                elif tag in ("td", "th") and self._cur_row is not None:
                    cell_text = "".join(self._cur_cell).strip().replace("\n", " ")
                    self._cur_row.append(cell_text)
                    self._in_cell = False

            def handle_data(self, data):
                if self._in_cell:
                    self._cur_cell.append(data)

        parser = _TableParser()
        parser.feed(html)
        if not parser.rows:
            return html

        # 转为 Markdown 表格
        lines: list[str] = []
        for i, row in enumerate(parser.rows):
            if not row:
                continue
            lines.append("| " + " | ".join(row) + " |")
            if i == 0:
                lines.append("| " + " | ".join("---" for _ in row) + " |")
        return "\n".join(lines) if lines else html
    except Exception:
        return html


def _ocr_image_with_layout(image) -> str:
    """用 PaddleX layout_parsing 做版面感知 OCR，返回结构化文本。

    相比 _ocr_image（纯 OCR 拍平）的改进：
    - 识别 doc_title / paragraph_title / text / table 等区域类型
    - paragraph_title 转为 `## 标题`，让 chunker 按标题切分
    - table 区域的 HTML 转 Markdown 表格，保留行列结构
    - 丢弃 header / footer 噪音

    Args:
        image: PIL.Image 对象

    Returns:
        结构化文本（带 Markdown 标记），或空字符串（不可用/失败时）
    """
    pipeline = _get_layout_pipeline()
    if pipeline is None:
        return ""

    import numpy as np  # type: ignore
    import io

    # PIL → 图片路径或 numpy（PaddleX predict 接受路径或 ndarray）
    arr = np.array(image)
    try:
        result = pipeline.predict(input=arr)
    except Exception:
        # ndarray 输入失败时尝试保存临时文件再传路径
        import tempfile
        tmp_path: Optional[str] = None
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            image.save(tmp, format="PNG")
            tmp_path = tmp.name
        try:
            result = pipeline.predict(input=tmp_path)
        finally:
            if tmp_path:
                _os_unlink(tmp_path)

    parts: list[str] = []
    for r in result:
        res = r.json.get("res", {}) if hasattr(r, "json") else {}
        parsing_list = res.get("parsing_res_list", [])
        for block in parsing_list:
            label = block.get("block_label", "")
            content = block.get("block_content", "")
            if not content or not content.strip():
                continue
            # 按区域类型结构化
            if label == "doc_title":
                parts.append("# " + content.strip())
            elif label == "paragraph_title":
                parts.append("## " + content.strip())
            elif label in ("text", "other_text"):
                parts.append(content.strip())
            elif label == "table":
                # 表格 HTML → Markdown
                md_table = _html_table_to_markdown(content)
                parts.append(md_table)
            elif label in ("header", "footer"):
                # 丢弃页眉页脚噪音
                continue
            elif label == "figure":
                # 图片区域跳过（无文本）
                continue
            else:
                # 其他类型（印章等）保留文本
                parts.append(content.strip())
        break  # 只取第一页/第一张图

    return _clean_ocr_text("\n\n".join(parts)) if parts else ""


def _os_unlink(path: str) -> None:
    """安全删除临时文件。"""
    try:
        import os
        os.unlink(path)
    except Exception:
        pass


def _ocr_image(image) -> str:
    """对 PIL Image 做 OCR，返回识别文本。

    优先级：
    1. PaddleX layout_parsing（版面感知，保留标题/表格结构，需 enable_layout_parsing=True）
    2. PaddleOCR predict（纯 OCR 文本，精度高）
    3. Tesseract（降级方案，外部预处理）

    所有 OCR 输出统一走 _clean_ocr_text 后处理。

    Args:
        image: PIL.Image 对象

    Returns:
        识别出的文本（已清理）
    """
    from config import settings

    # 1. 优先 PaddleX layout_parsing（版面感知，保留结构）
    if settings.enable_layout_parsing:
        try:
            layout_text = _ocr_image_with_layout(image)
            if layout_text.strip():
                return layout_text  # _ocr_image_with_layout 内部已调用 _clean_ocr_text
        except Exception:
            # layout_parsing 失败时降级到 PaddleOCR predict（某些图片版面分析兼容性问题）
            pass

    # 2. 降级 PaddleOCR predict（纯 OCR 文本）
    paddle_text = _ocr_image_paddle(image)
    if paddle_text.strip():
        return _clean_ocr_text(paddle_text)

    # 3. 再降级 Tesseract（外部预处理提升识别率）
    try:
        processed = _preprocess_image(image)
        tesseract_text = _ocr_image_tesseract(processed)
        if tesseract_text.strip():
            return _clean_ocr_text(tesseract_text)
        return ""
    except Exception:
        return ""


def _ocr_pdf_page(page) -> str:
    """对 PDF 单页做 OCR（用于扫描版 PDF）。

    Args:
        page: fitz.Page 对象

    Returns:
        识别出的文本
    """
    import fitz  # PyMuPDF
    # 渲染页面为图片（DPI 200，PaddleOCR 内部会做超分辨率）
    pix = page.get_pixmap(dpi=200)
    from PIL import Image  # type: ignore
    import io
    image = Image.open(io.BytesIO(pix.tobytes("png")))
    return _ocr_image(image)


def _render_pdf_page_to_png(page) -> bytes:
    """渲染 PDF 单页为 PNG 字节流（必须在主线程调用，PyMuPDF 非线程安全）。

    Args:
        page: fitz.Page 对象

    Returns:
        PNG 字节流
    """
    import fitz  # PyMuPDF
    pix = page.get_pixmap(dpi=200)
    return pix.tobytes("png")


def _ocr_png_bytes(png_bytes: bytes) -> str:
    """对 PNG 字节流做 OCR（线程安全，可在工作线程中调用）。

    与 _ocr_pdf_page 的区别：跳过 PyMuPDF 渲染步骤，
    直接从 PNG 字节流加载 PIL Image 后调用 _ocr_image。
    用于并行 OCR 场景：主线程预渲染所有页面，工作线程并行 OCR。

    Args:
        png_bytes: PNG 图片字节流

    Returns:
        识别出的文本
    """
    from PIL import Image  # type: ignore
    import io
    image = Image.open(io.BytesIO(png_bytes))
    return _ocr_image(image)


def _parallel_ocr_pages(
    page_indices: list[int],
    png_bytes_list: list[bytes],
    max_workers: int | None = None,
) -> dict[int, str]:
    """并行对多页 PNG 字节流做 OCR。

    策略：
    - PaddleOCR 调用通过 _paddle_ocr_lock 串行化（线程安全）
    - Tesseract 走子进程，可真正并行
    - 实际加速来源：图片预处理（numpy/PIL）+ Tesseract 子进程并行

    Args:
        page_indices: 页码列表（1-based）
        png_bytes_list: 对应的 PNG 字节流列表
        max_workers: 并发度（默认从 IMA_PDF_OCR_WORKERS 读取，否则 2）

    Returns:
        {page_index: ocr_text} 映射，失败的页码值为空字符串
    """
    import os
    from concurrent.futures import ThreadPoolExecutor, as_completed

    if not page_indices:
        return {}

    if max_workers is None:
        try:
            max_workers = int(os.environ.get("IMA_PDF_OCR_WORKERS", "2"))
        except (TypeError, ValueError):
            max_workers = 2
    # 并发度不超过任务数
    max_workers = max(1, min(max_workers, len(page_indices)))

    results: dict[int, str] = {}
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_idx = {
            executor.submit(_ocr_png_bytes, png): idx
            for idx, png in zip(page_indices, png_bytes_list)
        }
        for future in as_completed(future_to_idx):
            idx = future_to_idx[future]
            try:
                results[idx] = future.result()
            except Exception:
                results[idx] = ""
    return results


# ============================================================
# 各格式解析函数
# ============================================================

def _extract_pdf_page_blocks(page) -> list[dict]:
    """用 get_text("dict") 提取单页文本块，按阅读顺序排序。

    相比 page.get_text() 的改进：
    - 按 (y0, x0) 排序，恢复阅读顺序（多栏布局下避免左右栏交错）
    - 返回块列表，每块是 span 文本拼接，保留段落边界
    - 检测多栏布局，按栏内顺序输出
    - 保留 font_size 供下游推断标题层级

    Args:
        page: fitz.Page 对象

    Returns:
        该页的文本块列表，每块是 dict：
        {"text": str, "font_size": float, "x0/y0/x1/y1": float}
        （已按阅读顺序排序）
    """
    import fitz  # noqa: F401  PyMuPDF
    page_dict = page.get_text("dict")
    blocks = page_dict.get("blocks", [])

    # 提取文本块（跳过图片块）
    text_blocks: list[dict] = []
    for b in blocks:
        if b.get("type", 0) != 0:  # 0=文本, 1=图片
            continue
        lines = b.get("lines", [])
        if not lines:
            continue
        # 拼接 block 内所有 line 的 span 文本
        block_text_parts: list[str] = []
        max_font_size = 0.0
        for line in lines:
            line_text = "".join(span.get("text", "") for span in line.get("spans", []))
            if line_text.strip():
                block_text_parts.append(line_text)
                for span in line.get("spans", []):
                    sz = span.get("size", 0)
                    if sz > max_font_size:
                        max_font_size = sz
        block_text = "\n".join(block_text_parts).strip()
        if not block_text:
            continue
        bbox = b.get("bbox", [0, 0, 0, 0])  # (x0, y0, x1, y1)
        text_blocks.append({
            "text": block_text,
            "x0": bbox[0], "y0": bbox[1], "x1": bbox[2], "y1": bbox[3],
            "font_size": max_font_size,
        })

    if not text_blocks:
        return []

    # ---- 多栏检测 ----
    # 若 block 的 x 坐标明显分成左右两组（gap > 页宽 15%），按栏排序
    page_width = page.rect.width if hasattr(page, "rect") else 0
    if page_width > 0 and len(text_blocks) >= 4:
        x0_values = sorted(b["x0"] for b in text_blocks)
        mid_x = page_width / 2
        left_count = sum(1 for x in x0_values if x < mid_x)
        right_count = len(x0_values) - left_count
        # 若左右各有至少 2 块，判定为双栏，按 (栏, y0, x0) 排序
        if left_count >= 2 and right_count >= 2:
            for b in text_blocks:
                b["col"] = 0 if b["x0"] < mid_x else 1
            text_blocks.sort(key=lambda b: (b["col"], b["y0"], b["x0"]))
        else:
            text_blocks.sort(key=lambda b: (b["y0"], b["x0"]))
    else:
        text_blocks.sort(key=lambda b: (b["y0"], b["x0"]))

    return text_blocks


def _infer_heading_level(block_font_size: float, body_font_size: float) -> int:
    """根据字号相对大小推断标题层级。

    策略：以正文字号为基准，相对大小判断层级
    - font_size >= body * 1.8 → H1 (文档标题级)
    - font_size >= body * 1.4 → H2 (章节级)
    - font_size >= body * 1.15 → H3 (小节级)
    - 其他 → 0 (正文，不标记)

    Args:
        block_font_size: 当前块的字号
        body_font_size: 该页正文字号（众数）

    Returns:
        标题层级 1/2/3，0 表示非标题
    """
    if body_font_size <= 0 or block_font_size <= 0:
        return 0
    ratio = block_font_size / body_font_size
    if ratio >= 1.8:
        return 1
    if ratio >= 1.4:
        return 2
    if ratio >= 1.15:
        return 3
    return 0


def _detect_body_font_size(blocks: list[dict]) -> float:
    """检测页面正文字号（出现次数最多的字号）。

    用众数作为正文字号基准，避免被标题字号干扰。
    """
    if not blocks:
        return 0.0
    from collections import Counter
    # 按块文本长度加权统计（正文块通常比标题块长）
    weighted: Counter = Counter()
    for b in blocks:
        sz = round(b.get("font_size", 0), 1)
        if sz > 0:
            # 用文本长度作权重，让正文字号占主导
            weighted[sz] += len(b.get("text", ""))
    if not weighted:
        return 0.0
    return weighted.most_common(1)[0][0]


def _mark_headings(blocks: list[dict]) -> None:
    """为块标注标题层级（原地修改，添加 'heading_level' 字段）。

    仅对短文本块（<= 80 字符）标记为标题，长文本即使字号大也不标
    （可能是大字号的引言段落）。
    """
    body_size = _detect_body_font_size(blocks)
    if body_size <= 0:
        for b in blocks:
            b["heading_level"] = 0
        return
    for b in blocks:
        text = b.get("text", "").strip()
        # 标题判定：字号达标 + 文本短 + 不以句号结尾
        if (
            len(text) <= 80
            and not text.endswith(("。", "！", "？", "；", "."))
            and _infer_heading_level(b.get("font_size", 0), body_size) > 0
        ):
            b["heading_level"] = _infer_heading_level(b["font_size"], body_size)
        else:
            b["heading_level"] = 0


def _bbox_in_table(bbox: list, table_bbox: tuple, tolerance: float = 3.0) -> bool:
    """判断文本块 bbox 是否在表格 bbox 内（带容差）。

    Args:
        bbox: 文本块 bbox [x0, y0, x1, y1]
        table_bbox: 表格 bbox (x0, y0, x1, y1)
        tolerance: 容差（点），允许块稍微超出表格边界
    """
    if not bbox or not table_bbox:
        return False
    x0, y0, x1, y1 = bbox[0], bbox[1], bbox[2], bbox[3]
    tx0, ty0, tx1, ty1 = table_bbox
    return (
        x0 >= tx0 - tolerance and y0 >= ty0 - tolerance
        and x1 <= tx1 + tolerance and y1 <= ty1 + tolerance
    )


def _table_to_markdown(table_data: list, max_rows: int = 30) -> str:
    """把 find_tables().extract() 的二维数组转成 Markdown 表格。

    Args:
        table_data: 二维列表，每行是单元格文本列表
        max_rows: 最大行数（防止超大表格爆 token）

    Returns:
        Markdown 表格字符串，如：
        | 列1 | 列2 |
        |---|---|
        | 值1 | 值2 |
    """
    if not table_data:
        return ""

    # 截断超长表格
    if len(table_data) > max_rows:
        table_data = table_data[:max_rows]
        truncated = True
    else:
        truncated = False

    # 清理单元格文本（去换行、去首尾空格）
    def _clean(cell) -> str:
        if cell is None:
            return ""
        return str(cell).replace("\n", " ").replace("|", "\\|").strip()

    rows = [[_clean(c) for c in row] for row in table_data]
    if not rows:
        return ""

    # 确定列数（取最大）
    col_count = max(len(row) for row in rows)
    if col_count == 0:
        return ""

    # 补齐每行列数
    for row in rows:
        while len(row) < col_count:
            row.append("")

    # 第一行作表头，如果只有一行则补空表头
    if len(rows) == 1:
        header = rows[0]
        body = []
    else:
        header = rows[0]
        body = rows[1:]

    lines = []
    # 表头
    lines.append("| " + " | ".join(header) + " |")
    # 分隔行
    lines.append("|" + "|".join(["---"] * col_count) + "|")
    # 数据行
    for row in body:
        lines.append("| " + " | ".join(row) + " |")

    if truncated:
        lines.append(f"*（表格已截断，共显示 {max_rows} 行）*")

    return "\n".join(lines)


def _mark_table_blocks(page, blocks: list[dict]) -> list[tuple]:
    """用 find_tables() 检测表格，转 Markdown 表格插入 blocks。

    策略：
    1. 用 page.find_tables() 找到所有表格
    2. 把表格区域的文本块从 blocks 中移除（避免重复）
    3. 把 Markdown 表格作为新块按 y 坐标位置插入 blocks

    原地修改 blocks（移除表格内文本块，插入表格块）。
    返回表格 bbox 列表（供调试/统计）。
    """
    try:
        tables = page.find_tables()
    except Exception:
        return []

    table_list = getattr(tables, "tables", []) or []
    if not table_list:
        return []

    table_bboxes: list[tuple] = []
    table_dicts: list[dict] = []

    for tbl in table_list:
        try:
            data = tbl.extract()
        except Exception:
            continue
        md = _table_to_markdown(data)
        if not md.strip():
            continue
        bbox = tuple(tbl.bbox) if tbl.bbox else None
        table_bboxes.append(bbox)
        # 用表格中心点 y 作排序键
        y_center = (bbox[1] + bbox[3]) / 2 if bbox else 0
        table_dicts.append({
            "text": md,
            "font_size": 0,  # 表格块不参与标题判定
            "x0": bbox[0] if bbox else 0,
            "y0": bbox[1] if bbox else 0,
            "x1": bbox[2] if bbox else 0,
            "y1": bbox[3] if bbox else 0,
            "heading_level": 0,
            "is_table": True,
        })

    if not table_dicts:
        return []

    # 移除在表格 bbox 内的文本块
    filtered_blocks = []
    for b in blocks:
        in_table = False
        for tbbox in table_bboxes:
            if _bbox_in_table([b["x0"], b["y0"], b["x1"], b["y1"]], tbbox):
                in_table = True
                break
        if not in_table:
            filtered_blocks.append(b)

    # 清空原 blocks 再重新填充
    blocks.clear()
    blocks.extend(filtered_blocks)

    # 把表格块按 y 坐标插入到正确位置
    for td in table_dicts:
        blocks.append(td)

    # 重新排序
    blocks.sort(key=lambda b: (b.get("y0", 0), b.get("x0", 0)))

    return table_bboxes


def _detect_header_footer(page_texts: list[list[str]], min_repeat: int = 3) -> tuple[set[str], set[str]]:
    """检测 PDF 的页眉页脚（在多页重复出现的短文本块）。

    Args:
        page_texts: 每页的文本块列表
        min_repeat: 至少在 N 页出现才算页眉页脚

    Returns:
        (header_set, footer_set) 需要过滤的文本集合
    """
    if len(page_texts) < min_repeat:
        return set(), set()

    from collections import Counter
    # 统计每页首块和末块的出现频次
    first_blocks: Counter = Counter()
    last_blocks: Counter = Counter()
    for blocks in page_texts:
        if blocks:
            # 归一化：去除空白和页码数字
            first = re.sub(r"\s+", "", blocks[0])[:50]
            last = re.sub(r"\s+", "", blocks[-1])[:50]
            if first:
                first_blocks[first] += 1
            if last:
                last_blocks[last] += 1

    # 出现在 >= min_repeat 页的首/末块视为页眉页脚
    page_count = len(page_texts)
    threshold = max(min_repeat, page_count // 2)
    header_set = {t for t, c in first_blocks.items() if c >= threshold}
    footer_set = {t for t, c in last_blocks.items() if c >= threshold}

    # 额外检测：纯数字（页码）或纯页码+短文字
    for blocks in page_texts:
        if blocks:
            for idx in (0, -1):
                txt = re.sub(r"\s+", "", blocks[idx])[:50]
                # 纯数字 或 "数字 文字" 形式的页码
                if re.match(r"^\d{1,3}$", txt) or re.match(r"^.{0,10}\d{1,3}$", txt):
                    if idx == 0:
                        header_set.add(txt)
                    else:
                        footer_set.add(txt)

    return header_set, footer_set


def _is_header_footer(text: str, header_set: set, footer_set: set) -> bool:
    """判断文本块是否为页眉页脚。"""
    normalized = re.sub(r"\s+", "", text)[:50]
    if not normalized:
        return False
    if normalized in header_set or normalized in footer_set:
        return True
    # 模糊匹配：页眉页脚集合中的项是否是当前文本的前缀/后缀
    # 限制：匹配项长度 ≥ 4 字符，且原文长度不超过匹配项 + 5 字符
    # （避免"附件"误匹配"附件1：XXX"等正文标题）
    for h in header_set:
        if len(h) >= 4 and len(normalized) <= len(h) + 5 and (normalized.startswith(h) or normalized.endswith(h)):
            return True
    for f in footer_set:
        if len(f) >= 4 and len(normalized) <= len(f) + 5 and (normalized.startswith(f) or normalized.endswith(f)):
            return True
    return False


def _merge_visual_line_breaks(text: str) -> str:
    """合并 PDF 视觉换行（单 \\n），保留段落分隔（双 \\n）。

    PDF 提取的文本中，单 \\n 是视觉换行（同一行的折行），双 \\n 是段落分隔。
    对单 \\n 做合并处理：
    - 行尾是中文标点（。！？；）或全角符号 → 保留换行（段落结束）
    - 行尾是英文句号/问号/感叹号+空格 → 保留换行
    - 行首是中文标点（、，。等）→ 合并到上一行（标点不该行首）
    - 行首是列表标记（1. （1） • 等）→ 保留换行
    - 其他情况：合并（去掉 \\n，中文直接拼接，英文加空格）

    用于 _parse_pdf 的文本后处理，让 BM25 bigram 短语匹配恢复有效。
    """
    if not text:
        return text

    # 先规范化 Windows 换行
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # 把 3+ 换行压缩为 2 个（段落分隔）
    text = re.sub(r"\n{3,}", "\n\n", text)

    lines = text.split("\n")
    result: list[str] = []
    buffer = ""

    # 中文行尾标点（段落结束标志）
    cn_end_puncts = set("。！？；：…」』）】》")
    # 中文行首标点（不该行首，应合并到上一行）
    cn_start_puncts = set("、，。！？；：…」』）】》")
    # 列表/标题标记（行首应保留换行）
    list_pattern = re.compile(
        r"^(?:"
        r"\d+[.、）)]"          # 1. 1、 1） 1)
        r"|[（(]\d+[）)]"       # （1） (1)
        r"|[一二三四五六七八九十]+[、.）)]"  # 一、 二.
        r"|第[一二三四五六七八九十百千]+[章节条]"  # 第一章
        r"|●|•|◆|■|★|☆"      # 项目符号
        r"|[【\[《「『]"       # 引号开头
        r"|#{1,6}\s"           # Markdown 标题
        r")"
    )

    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            # 空行：如果 buffer 非空，可能是段落分隔
            if buffer:
                result.append(buffer)
                buffer = ""
            result.append("")
            continue

        if not buffer:
            buffer = stripped
            continue

        buffer_end = buffer[-1] if buffer else ""
        line_start = stripped[0] if stripped else ""

        # 判断是否应该保留换行（新起一段）
        should_new_line = False

        # 1. 上一行以段落结束标点结尾
        if buffer_end in cn_end_puncts:
            should_new_line = True
        # 2. 上一行以英文句末标点结尾
        elif buffer_end in {".", "!", "?", ";"} and len(buffer) >= 2 and (buffer[-2] == " " or buffer[-2].isascii()):
            should_new_line = True
        # 3. 当前行以列表/标题标记开头
        elif list_pattern.match(stripped):
            should_new_line = True
        # 4. 当前行首是中文标点（不该行首）→ 合并
        elif line_start in cn_start_puncts:
            should_new_line = False
        # 5. 当前行是全大写英文（可能是标题）
        elif stripped.isupper() and len(stripped) > 3 and stripped.isascii():
            should_new_line = True
        # 6. buffer 和当前行长度都较短（可能是标题行）
        elif len(buffer) < 20 and len(stripped) < 20 and not cn_start_puncts.intersection(stripped):
            # 短行可能是标题，保留换行
            should_new_line = True

        if should_new_line:
            result.append(buffer)
            buffer = stripped
        else:
            # 合并：中文直接拼接，中英/英英之间加空格
            if buffer_end.isascii() and line_start.isascii():
                buffer = buffer + " " + stripped
            elif buffer_end.isascii() or line_start.isascii():
                buffer = buffer + " " + stripped
            else:
                buffer = buffer + stripped

    if buffer:
        result.append(buffer)

    return "\n".join(result).strip()


def _parse_pdf(file_path: Path) -> ParsedDocument:
    """解析 PDF：先尝试提取文本层（结构感知），若为空则走 OCR。

    改进点：
    - 用 get_text("dict") 按阅读顺序提取，检测多栏布局
    - 检测并去除重复的页眉页脚
    - 合并视觉换行（单 \\n），保留段落分隔（双 \\n）
    - 保留页码分隔标记（\\n\\n--- Page N ---\\n\\n），供下游分块使用
    - O3: 用 font_size 推断标题层级，转成 Markdown #/##/### 标记
    - O1: 用 find_tables() 还原表格结构为 Markdown 表格
    - O2: 文本块少的页提取图片块走 OCR
    - P1-性能: 扫描页并行 OCR（PyMuPDF 渲染保持串行，OCR 调用并行）

    Args:
        file_path: PDF 文件路径

    Returns:
        ParsedDocument
    """
    import fitz  # PyMuPDF

    # 每页的块列表（dict 格式，含 font_size/heading_level）
    raw_page_blocks: list[list[dict]] = []
    page_count = 0
    ocr_pages: list[int] = []
    ocr_failed_pages: list[int] = []

    # Phase 1: 串行扫描所有页面，提取文本块；需要 OCR 的页面预渲染为 PNG 字节流
    # PyMuPDF 非线程安全，渲染必须在主线程串行执行
    pending_ocr: list[tuple[int, bytes]] = []  # [(page_index_1based, png_bytes)]
    pending_ocr_slots: list[int] = []  # 在 raw_page_blocks 中的位置（用于回填）

    with fitz.open(file_path) as doc:
        page_count = len(doc)
        for i, page in enumerate(doc):
            blocks = _extract_pdf_page_blocks(page)
            # O1: 提取表格并标记表格区域，避免文本块重复
            table_bboxes = _mark_table_blocks(page, blocks)
            # O3: 标注标题层级
            _mark_headings(blocks)
            total_text = "\n".join(b["text"] for b in blocks).strip()

            # 判断该页是否需要 OCR：文本层几乎为空（< 50 字符）
            if total_text and len(total_text) >= 50:
                raw_page_blocks.append(blocks)
            else:
                # 需 OCR：先在主线程渲染为 PNG（PyMuPDF 非线程安全）
                slot = len(raw_page_blocks)
                raw_page_blocks.append([])  # 占位，后续回填
                if _check_ocr():
                    try:
                        png_bytes = _render_pdf_page_to_png(page)
                        pending_ocr.append((i + 1, png_bytes))
                        pending_ocr_slots.append(slot)
                    except Exception:
                        ocr_failed_pages.append(i + 1)
                else:
                    # OCR 不可用，直接标记为失败
                    ocr_failed_pages.append(i + 1)

    # Phase 2: 并行 OCR 所有待处理页面
    if pending_ocr:
        ocr_results = _parallel_ocr_pages(
            [idx for idx, _ in pending_ocr],
            [png for _, png in pending_ocr],
        )
        # 回填结果
        for (page_idx, _), slot in zip(pending_ocr, pending_ocr_slots):
            ocr_text = ocr_results.get(page_idx, "")
            if ocr_text.strip():
                raw_page_blocks[slot] = [{
                    "text": ocr_text, "font_size": 0, "heading_level": 0,
                }]
                ocr_pages.append(page_idx)
            else:
                ocr_failed_pages.append(page_idx)

    # 检测页眉页脚（用 text 字符串）
    header_set, footer_set = _detect_header_footer(
        [[b["text"] for b in blocks] for blocks in raw_page_blocks]
    )

    # 组装最终文本：去页眉页脚 + 标题转 # + 合并视觉换行 + 页码分隔标记
    page_sections: list[str] = []
    for i, blocks in enumerate(raw_page_blocks):
        if not blocks:
            continue
        filtered: list[str] = []
        for b in blocks:
            if not _is_header_footer(b["text"], header_set, footer_set):
                # O3: 标题块转 Markdown 标记
                level = b.get("heading_level", 0)
                if level > 0:
                    filtered.append("#" * level + " " + b["text"])
                else:
                    filtered.append(b["text"])
        if not filtered:
            continue
        # O1: 表格块（含 |---|）不参与视觉换行合并，用双换行与正文分隔
        # 策略：把表格块用占位符替换 → 合并正文换行 → 换回表格
        _TABLE_PLACEHOLDER_PREFIX = "\x00TABLE_"
        placeholders: dict[str, str] = {}
        for idx_b, b in enumerate(filtered):
            if "|---|" in b:
                key = f"{_TABLE_PLACEHOLDER_PREFIX}{idx_b}\x00"
                placeholders[key] = b
                filtered[idx_b] = key
        page_text = "\n".join(filtered)
        # 合并视觉换行（只作用于正文，表格已替换为占位符）
        page_text = _merge_visual_line_breaks(page_text)
        # 换回表格块
        for key, original in placeholders.items():
            page_text = page_text.replace(key, original)
        # 加页码分隔标记（供 chunker 识别页边界）
        page_sections.append(f"--- Page {i + 1} ---\n{page_text}")

    text = "\n\n".join(page_sections).strip()

    meta: Dict[str, str] = {"page_count": str(page_count)}
    if ocr_pages:
        meta["ocr_pages"] = ",".join(str(p) for p in ocr_pages)
        meta["ocr_used"] = "true"
    if ocr_failed_pages:
        meta["ocr_failed_pages"] = ",".join(str(p) for p in ocr_failed_pages)
    if header_set:
        meta["header_removed"] = "true"
    if footer_set:
        meta["footer_removed"] = "true"

    return ParsedDocument(
        text=text,
        title=file_path.stem,
        file_path=file_path,
        file_type=".pdf",
        language="unknown",
        meta=meta,
        format_tag="pdf",
    )


def _parse_image(file_path: Path) -> ParsedDocument:
    """解析图片：走 OCR 提取文本。

    Args:
        file_path: 图片文件路径

    Returns:
        ParsedDocument
    """
    if not _check_ocr():
        # OCR 不可用，返回空文档（让上层跳过）
        return ParsedDocument(
            text="",
            title=file_path.stem,
            file_path=file_path,
            file_type=file_path.suffix.lower(),
            meta={"ocr_unavailable": "true"},
            format_tag="image",
        )

    from PIL import Image  # type: ignore

    try:
        # Bug 10 修复：用 with 语句确保文件句柄关闭
        # 批量 OCR 场景下未关闭句柄会触发 ResourceWarning 并可能耗尽 fd
        with Image.open(file_path) as image:
            text = _ocr_image(image)
    except Exception as e:
        raise ParseError(f"图片 OCR 失败 [{file_path.name}]: {type(e).__name__}: {e}")

    return ParsedDocument(
        text=text.strip(),
        title=file_path.stem,
        file_path=file_path,
        file_type=file_path.suffix.lower(),
        language="unknown",
        meta={"ocr_used": "true"},
        format_tag="image",
    )


def _parse_docx(file_path: Path) -> ParsedDocument:
    """解析 Word .docx：提取段落 + 表格内容，保留 Heading 样式层级。

    样式感知：
    - 识别 `paragraph.style.name`（Heading 1/2/3/Title 等），转为 Markdown 风格的 `#`/`##`/`###` 标记
    - 列表样式（List Bullet/List Number）保留为 `-`/`1.` 前缀
    - 让下游 chunker 能按标题层级切分，并填充 heading 字段供 reranker 使用
    """
    from docx import Document

    doc = Document(str(file_path))
    parts: list[str] = []

    # 1. 段落文本（样式感知）
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style_name = (p.style.name or "").strip() if p.style else ""
        # Heading 1/2/3/4/5/6 → Markdown #/##/###/####/#####/######
        if style_name.lower().startswith("heading"):
            try:
                level = int(style_name.split()[-1])
                level = max(1, min(6, level))
                parts.append("#" * level + " " + text)
            except (ValueError, IndexError):
                parts.append("# " + text)  # 无法解析层级，按 H1 处理
        elif style_name.lower() == "title":
            parts.append("# " + text)
        elif style_name.lower().startswith("list bullet"):
            parts.append("- " + text)
        elif style_name.lower().startswith("list number"):
            parts.append("1. " + text)
        else:
            parts.append(text)

    # 2. 表格内容（Markdown 表格格式，与 PDF 统一）
    for tbl_idx, table in enumerate(doc.tables):
        rows = table.rows
        if not rows:
            continue
        # 提取二维列表
        table_data = []
        for row in rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        # 转成 Markdown 表格
        md_table = _table_to_markdown(table_data)
        if md_table:
            parts.append(md_table)

    text = "\n\n".join(parts)
    return ParsedDocument(
        text=text,
        title=file_path.stem,
        file_path=file_path,
        file_type=".docx",
        meta={"paragraph_count": str(len(doc.paragraphs))},
        format_tag="docx",
    )


def _find_libreoffice() -> str | None:
    """查找 LibreOffice 可执行文件（跨平台）。

    查找顺序：
    1. PATH 中的 soffice / libreoffice
    2. macOS 应用包 /Applications/LibreOffice.app/Contents/MacOS/soffice

    Returns:
        可执行文件路径，未找到返回 None
    """
    from shutil import which

    # 1. PATH 查找
    for cmd in ("soffice", "libreoffice"):
        path = which(cmd)
        if path:
            return path

    # 2. macOS 应用包
    if sys.platform == "darwin":
        macos_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if Path(macos_path).exists():
            return macos_path

    return None


def _doc_to_docx_via_libreoffice(file_path: Path) -> Path | None:
    """用 LibreOffice headless 把 .doc 转成 .docx。

    Args:
        file_path: .doc 文件路径

    Returns:
        转换后的 .docx 临时文件路径，失败返回 None
    """
    import subprocess
    import tempfile

    soffice = _find_libreoffice()
    if not soffice:
        return None

    # 用临时目录接收转换结果
    out_dir = tempfile.mkdtemp(prefix="ima_doc_convert_")
    try:
        result = subprocess.run(
            [
                soffice, "--headless",
                "--convert-to", "docx",
                "--outdir", out_dir,
                str(file_path),
            ],
            capture_output=True,
            text=True,
            timeout=60,  # 防止卡死
        )
        if result.returncode != 0:
            return None
        # 转换后的文件名：原文件名 + .docx
        docx_path = Path(out_dir) / (file_path.stem + ".docx")
        if docx_path.exists():
            return docx_path
        return None
    except (subprocess.TimeoutExpired, Exception):
        return None


def _parse_doc(file_path: Path) -> ParsedDocument:
    """解析旧版 Word .doc：优先用 LibreOffice 转 DOCX，降级用 textutil。

    转换策略：
    1. 优先用 LibreOffice headless 转成 .docx → 复用 _parse_docx（保留样式 + 表格）
    2. LibreOffice 不可用时，macOS 降级用 textutil 转 txt（丢样式）
    3. 都不可用则报错

    Args:
        file_path: .doc 文件路径

    Returns:
        ParsedDocument

    Raises:
        ParseError: 无可用转换工具或转换失败
    """
    import subprocess
    from shutil import which

    # 1. 优先用 LibreOffice 转 DOCX（保留样式 + 表格）
    docx_path = _doc_to_docx_via_libreoffice(file_path)
    if docx_path is not None:
        try:
            parsed = _parse_docx(docx_path)
            # 修正 file_path 和 file_type 为原始 .doc
            parsed.file_path = file_path
            parsed.file_type = ".doc"
            parsed.meta["converter"] = "libreoffice"
            return parsed
        finally:
            # 清理临时文件
            try:
                docx_path.unlink()
                docx_path.parent.rmdir()
            except Exception:
                pass

    # 2. 降级：macOS textutil 转 txt（丢样式）
    if which("textutil") is not None:
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", str(file_path)],
                capture_output=True,
                text=True,
                check=True,
            )
            text = result.stdout.strip()
            return ParsedDocument(
                text=text,
                title=file_path.stem,
                file_path=file_path,
                file_type=".doc",
                meta={"converter": "textutil"},
                format_tag="docx",
            )
        except subprocess.CalledProcessError as e:
            raise ParseError(f"textutil 转换失败 [{file_path.name}]: {e.stderr}") from e

    # 3. 都不可用
    raise ParseError(
        f"解析 .doc 需要安装 LibreOffice（推荐）或在 macOS 上使用 textutil: {file_path.name}"
    )


def _parse_xlsx(file_path: Path) -> ParsedDocument:
    """解析 Excel .xlsx：每行序列化为「表头=值」的键值对文本。

    优化点（解决行被切断、表头数据关联丢失问题）：
    1. 每行数据前缀表头字段名，形成 "列名1=值1 | 列名2=值2" 的语义完整单元
    2. 每行作为一个独立记录，永不切断单行
    3. 表头信息重复到每行，向量嵌入能感知字段语义
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    sheet_texts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # 第一行作为表头
        header_row = rows[0]
        headers = [
            (str(c).strip() if c is not None else f"列{i+1}")
            for i, c in enumerate(header_row)
        ]

        # 数据行序列化为键值对
        # 每行前缀「文档=DocTitle」注入文档标题语义，让 BM25 能命中标题里的关键词
        # 例：xlsx 标题"2025年智能服务终端配置安装情况"含"智能"二字，
        # 但原 chunk 内容只有 sheet 名"钱塘安装情况"，缺"智能"
        # 注入后 BM25 能命中"智能终端怎么安装"类查询
        doc_title = file_path.stem
        records: list[str] = []
        for row_idx, row in enumerate(rows[1:], start=2):
            if all(cell is None or str(cell).strip() == "" for cell in row):
                continue
            pairs = []
            for i, cell in enumerate(row):
                if i >= len(headers):
                    break
                val = "" if cell is None else str(cell).strip()
                # 跳过完全空值的列
                if val:
                    pairs.append(f"{headers[i]}={val}")
            if pairs:
                records.append(f"[{sheet_name}] 文档={doc_title} | " + " | ".join(pairs))

        if records:
            sheet_texts.append("\n".join(records))

    wb.close()

    text = "\n\n".join(sheet_texts)
    return ParsedDocument(
        text=text,
        title=file_path.stem,
        file_path=file_path,
        file_type=".xlsx",
        meta={"sheet_count": str(len(wb.sheetnames))},
        format_tag="excel",
    )


def _parse_pptx(file_path: Path) -> ParsedDocument:
    """解析 PowerPoint .pptx：提取每页幻灯片文本 + 表格 + 备注。

    优化：
    - 提取 slide 标题作为 metadata（供 reranker 和引用溯源使用）
    - 提取 notes_slide 备注（讲者备注，常含关键补充信息）
    - 在 `## Slide N` 后注入 `[标题] xxx` 标记，保留标题语义
    """
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slide_texts: list[str] = []
    slide_titles: list[str] = []  # 用于 meta
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        slide_title = ""

        # 1. 提取 slide 标题（优先用 placeholders.title）
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                slide_title = slide.shapes.title.text.strip()
        except Exception:
            pass

        # 2. 普通文本框 + 表格
        for shape in slide.shapes:
            # 跳过 title shape（已单独提取）
            if slide.shapes.title is not None and shape == slide.shapes.title:
                continue
            # 2.1 普通文本框
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            # 2.2 表格 shape（Markdown 表格格式，与 PDF/DOCX 统一）
            if shape.has_table:
                tbl = shape.table
                rows = tbl.rows
                if not rows:
                    continue
                # 提取二维列表
                table_data = []
                for row in rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                # 转成 Markdown 表格
                md_table = _table_to_markdown(table_data)
                if md_table:
                    texts.append(md_table)

        # 3. 备注（notes_slide）— 讲者备注常含关键补充信息
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    texts.append("[备注] " + notes_text)
        except Exception:
            pass

        if texts or slide_title:
            # 在 `## Slide N` 后注入标题标记（如有），保留标题语义供下游 chunker 使用
            header = f"## Slide {idx}"
            if slide_title:
                header += f"\n# {slide_title}"
            slide_texts.append(header + "\n" + "\n".join(texts))
            slide_titles.append(slide_title)

    text = "\n\n".join(slide_texts)
    return ParsedDocument(
        text=text,
        title=file_path.stem,
        file_path=file_path,
        file_type=".pptx",
        meta={
            "slide_count": str(len(prs.slides)),
            "slide_titles": " | ".join(t for t in slide_titles if t),
        },
        format_tag="ppt",
    )


def _html_to_markdown(html_text: str) -> str:
    """把 HTML 转成结构化 Markdown：保留 H1-H6 标题层级 + 表格。

    策略：
    1. 先用 trafilatura 提取正文（去导航/广告/侧栏）
    2. 用 trafilatura 的 extract 结果作为纯净 HTML 基础
    3. 但 trafilatura 2.0 的 markdown 输出实测不保留 # 标记
    4. 所以用 BeautifulSoup 直接解析原始 HTML，保留 H1-H6 + 表格 + 段落

    Args:
        html_text: 原始 HTML 字符串

    Returns:
        结构化 Markdown 文本（含 # 标题和 Markdown 表格）
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html_text, "html.parser")

    # 移除 script/style/nav/footer 等干扰元素
    for tag in soup.find_all(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    # 找到正文容器（优先 article/main，否则用 body）
    content = soup.find("article") or soup.find("main") or soup.body or soup

    parts: list[str] = []

    def _process_tag(element) -> None:
        """递归处理 HTML 元素，转成 Markdown。"""
        from bs4 import NavigableString, Tag

        if isinstance(element, NavigableString):
            text = str(element).strip()
            if text:
                parts.append(text)
            return

        if not isinstance(element, Tag):
            return

        name = element.name.lower()

        # 标题 H1-H6 → Markdown #/##/###
        if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(name[1])
            text = element.get_text(strip=True)
            if text:
                parts.append("#" * level + " " + text)
            return

        # 段落
        if name == "p":
            text = element.get_text(strip=True)
            if text:
                parts.append(text)
            return

        # 列表
        if name in ("ul", "ol"):
            for i, li in enumerate(element.find_all("li", recursive=False), 1):
                text = li.get_text(strip=True)
                if text:
                    prefix = f"{i}. " if name == "ol" else "- "
                    parts.append(prefix + text)
            return

        # 表格 → Markdown 表格（复用 _table_to_markdown）
        if name == "table":
            table_data: list[list[str]] = []
            for tr in element.find_all("tr"):
                row = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                if row:
                    table_data.append(row)
            if table_data:
                md = _table_to_markdown(table_data)
                if md:
                    parts.append(md)
            return

        # 引用块
        if name == "blockquote":
            text = element.get_text(strip=True)
            if text:
                parts.append("> " + text)
            return

        # 其他容器：递归处理子元素
        for child in element.children:
            _process_tag(child)

    _process_tag(content)

    # 合并并清理多余空行
    result = "\n\n".join(p for p in parts if p.strip())
    return result


def _parse_html(file_path: Path) -> ParsedDocument:
    """解析 HTML：trafilatura 提正文 + BeautifulSoup 保结构。

    两层处理：
    1. trafilatura.extract() 提取正文（去导航/广告/侧栏）
    2. _html_to_markdown() 保留 H1-H6 标题层级 + 表格结构

    如果 trafilatura 提取失败，降级为直接 BeautifulSoup 解析。
    """
    html_text = file_path.read_text(encoding="utf-8", errors="ignore")

    # 优先用 _html_to_markdown（保留结构）
    try:
        md_text = _html_to_markdown(html_text)
        if md_text and len(md_text.strip()) >= 20:
            return ParsedDocument(
                text=md_text.strip(),
                title=file_path.stem,
                file_path=file_path,
                file_type=".html",
                meta={"extractor": "bs4"},
                format_tag="html",
            )
    except Exception:
        pass

    # 降级：trafilatura 纯文本提取
    import trafilatura
    text = trafilatura.extract(html_text) or ""
    return ParsedDocument(
        text=text.strip(),
        title=file_path.stem,
        file_path=file_path,
        file_type=".html",
        meta={"extractor": "trafilatura"},
        format_tag="html",
    )


def _parse_plain(file_path: Path) -> ParsedDocument:
    """解析纯文本 / Markdown：直接读取。"""
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    ext = file_path.suffix.lower()
    # Markdown 文件打 markdown tag，其他打 text tag
    format_tag = "markdown" if ext in (".md", ".markdown") else "text"
    return ParsedDocument(
        text=text,
        title=file_path.stem,
        file_path=file_path,
        file_type=ext,
        format_tag=format_tag,
    )


def _parse_code(file_path: Path) -> ParsedDocument:
    """解析代码文件：保留原文 + 标注语言。"""
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    language = _CODE_LANGUAGES.get(file_path.suffix.lower(), "code")
    return ParsedDocument(
        text=text,
        title=file_path.stem,
        file_path=file_path,
        file_type=file_path.suffix.lower(),
        language=language,
        meta={"is_code": "true"},
    )


# ============================================================
# 解析器注册表
# ============================================================

_PARSER_MAP: Dict[str, Callable[[Path], ParsedDocument]] = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".doc": _parse_doc,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".html": _parse_html,
    ".htm": _parse_html,
    # 图片（OCR）
    ".png": _parse_image,
    ".jpg": _parse_image,
    ".jpeg": _parse_image,
    ".tif": _parse_image,
    ".tiff": _parse_image,
    ".bmp": _parse_image,
    ".webp": _parse_image,
}

# Markdown / 文本类
for ext in (".md", ".markdown", ".txt", ".log"):
    _PARSER_MAP[ext] = _parse_plain

# 注：代码文件按用户决策不入库，不再注册到 _PARSER_MAP


# ============================================================
# 公共 API
# ============================================================

def parse(file_path: Path | str) -> ParsedDocument:
    """解析单个文件，返回 ParsedDocument。

    Args:
        file_path: 文件路径（Path 或字符串）

    Returns:
        ParsedDocument

    Raises:
        FileNotFoundError: 文件不存在
        ParseError: 不支持的格式或解析失败
    """
    path = Path(file_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"文件不存在：{path}")
    if not path.is_file():
        raise ParseError(f"不是文件：{path}")

    ext = path.suffix.lower()
    parser = _PARSER_MAP.get(ext)
    if parser is None:
        raise ParseError(
            f"不支持的文件格式：{ext}。"
            f"支持的格式：{', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    try:
        return parser(path)
    except ParseError:
        raise
    except Exception as e:
        raise ParseError(f"解析失败 [{path.name}]: {type(e).__name__}: {e}") from e


def is_supported(file_path: Path | str) -> bool:
    """判断文件是否支持解析。"""
    return Path(file_path).suffix.lower() in SUPPORTED_EXTENSIONS
